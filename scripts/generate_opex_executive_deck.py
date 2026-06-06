#!/usr/bin/env python3
"""Generate a single-slide Deloitte-branded OpEx Intelligence executive summary."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DELOITTE_GREEN = RGBColor(0x86, 0xBC, 0x25)
DELOITTE_GREEN_DARK = RGBColor(0x43, 0xB0, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)
COOL_GRAY_11 = RGBColor(0x53, 0x56, 0x5A)
COOL_GRAY_6 = RGBColor(0xA7, 0xA9, 0xAC)
COOL_GRAY_1 = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUTPUT = Path(__file__).resolve().parents[1] / "deliverables" / "OpEx_Intelligence_Executive_Deck_Deloitte.pptx"
FONT = "Open Sans"


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _para(tf, text, *, size=10, bold=False, color=COOL_GRAY_11, italic=False, space_after=4):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
    if tf.paragraphs[0].text == "" and p is tf.paragraphs[0]:
        pass
    elif tf.paragraphs[0].text != "" or (len(tf.paragraphs) > 1):
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.space_after = Pt(space_after)
    return p


def _section_header(slide, left, top, width, title: str):
    hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.28))
    _set_fill(hdr, DELOITTE_GREEN)
    hdr.line.fill.background()
    box = _textbox(slide, left + Inches(0.08), top + Inches(0.03), width - Inches(0.12), Inches(0.22))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT


def _section_body(slide, left, top, width, height, bullets: list[str], size=8.5):
    box = _textbox(slide, left + Inches(0.06), top, width - Inches(0.1), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(size)
        p.font.color.rgb = COOL_GRAY_11
        p.font.name = FONT
        p.space_after = Pt(3)
        p.level = 0


def build_single_slide() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top brand bar
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    _set_fill(bar, DELOITTE_GREEN)
    bar.line.fill.background()

    # Deloitte wordmark
    brand = _textbox(slide, prs.slide_width - Inches(1.85), Inches(0.14), Inches(1.5), Inches(0.35))
    bp = brand.text_frame.paragraphs[0]
    bp.text = "Deloitte."
    bp.font.size = Pt(18)
    bp.font.bold = True
    bp.font.color.rgb = DELOITTE_GREEN
    bp.font.name = FONT
    bp.alignment = PP_ALIGN.RIGHT

    # Title block
    title = _textbox(slide, Inches(0.4), Inches(0.12), Inches(7.5), Inches(0.45))
    tp = title.text_frame.paragraphs[0]
    tp.text = "OpEx Intelligence Platform"
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = BLACK
    tp.font.name = FONT

    sub = _textbox(slide, Inches(0.4), Inches(0.52), Inches(8.5), Inches(0.3))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Agentic Cost Diagnostic for Indian Enterprises"
    sp.font.size = Pt(12)
    sp.font.color.rgb = DELOITTE_GREEN_DARK
    sp.font.name = FONT

    # 2×2 quadrant grid
    margin_l = Inches(0.4)
    gap = Inches(0.12)
    col_w = Inches(4.5)
    row_h = Inches(1.75)
    y1 = Inches(0.9)
    y2 = y1 + row_h + gap
    x1 = margin_l
    x2 = margin_l + col_w + gap

    sections = [
        (x1, y1, "WHAT IT IS", [
            "Consultant-deployable cost diagnostic asset deployed inside the client's environment",
            "Agentic engine: 36 skills + 11 sector packs — not SaaS, not ERP",
            "Identifies, sizes & stages 12–24 months of OpEx savings for Indian enterprises",
            "Built for Group CFO, Cost Transformation Office & advisory teams",
        ]),
        (x2, y1, "HOW IT WORKS", [
            "OBSERVE — ingest spend + context; classify data; strip PII",
            "PLAN — AI sequences skill DAG (profile → benchmark → model)",
            "ACT — parallel skills: peer bench, lever detect, savings model, evidence gather",
            "REFLECT — quality gates, evidence scoring, assumption register → board outputs",
        ]),
        (x1, y2, "EXPECTED OUTCOMES", [
            "2–4 hours from upload to insights · 12-week sprint to board commitment",
            "8–15% of addressable spend identified · 1.5–3% revenue capture over 18 months",
            "30–80 prioritized initiatives (P10/P50/P90, owner, payback, EBITDA bps)",
            "CFO brief · board deck · live cost room · MOR pack · PMO tracker",
        ]),
        (x2, y2, "WHY IT'S DIFFERENT", [
            "India-native: INR, GST, BRSR, DPDP, RBI/SEBI/IRDAI, CMIE/CRISIL benchmarks",
            "Consulting rigor at machine speed — collapses 16–24 weeks of analyst work",
            "C-suite narrative: bps EBITDA, ΔROCE, ΔEPS — audit-committee defensible",
            "Trust by design: client VPC, zero-retention LLM, immutable audit log, tear-down",
            "Evidence maturity scoring — every ₹ traces to source data & named assumptions",
        ]),
    ]

    for x, y, heading, bullets in sections:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, col_w, row_h
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = COOL_GRAY_1
        panel.line.color.rgb = DELOITTE_GREEN
        panel.line.width = Pt(0.75)

        _section_header(slide, x, y, col_w, heading)
        _section_body(slide, x, y + Inches(0.32), col_w, row_h - Inches(0.35), bullets)

    # Bottom pitch banner
    banner_y = Inches(4.55)
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), banner_y, Inches(9.2), Inches(0.55)
    )
    _set_fill(banner, DELOITTE_GREEN)
    banner.line.fill.background()

    pitch = _textbox(slide, Inches(0.55), banner_y + Inches(0.1), Inches(8.9), Inches(0.4))
    pp = pitch.text_frame.paragraphs[0]
    pp.text = (
        "Turns fragmented spend data into an audit-grade, board-ready savings portfolio — "
        "in hours, not months — built for how Indian enterprises commit cost transformation."
    )
    pp.font.size = Pt(11)
    pp.font.bold = True
    pp.font.color.rgb = WHITE
    pp.font.name = FONT
    pp.alignment = PP_ALIGN.CENTER

    # Footer
    footer_y = prs.slide_height - Inches(0.38)
    fline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), footer_y, Inches(9.2), Pt(1.5)
    )
    _set_fill(fline, COOL_GRAY_6)
    fline.line.fill.background()

    foot = _textbox(slide, Inches(0.4), footer_y + Inches(0.06), Inches(5), Inches(0.25))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Confidential — For discussion purposes only"
    fp.font.size = Pt(8)
    fp.font.color.rgb = COOL_GRAY_6
    fp.font.name = FONT

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_single_slide()
    print(f"Created: {path}")
